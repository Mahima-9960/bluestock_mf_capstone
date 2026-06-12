import os
from pptx import Presentation
from pptx.util import Inches, Pt

def create_bluestock_presentation():
    prs = Presentation()
    
    # Define slide layouts (0 = Title, 1 = Title + Content, 6 = Blank)
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    
    # --- SLIDE 1: Title ---
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "End-to-End Mutual Fund Analytics & Risk Engine"
    slide.placeholders[1].text = "Capstone Project Portfolio Milestone\nPresenter: Mahima Subhash Bhoite\nDate: June 2026"
    
    # --- SLIDE 2: Problem & Objectives ---
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Problem Statement & Project Objectives"
    tf = slide.placeholders[1].text_frame
    tf.text = "The Problem:"
    p = tf.add_paragraph()
    p.text = "• Financial analysts face data fragmentation when looking at historical NAV alongside investor behavior."
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Project Objectives:"
    p3 = tf.add_paragraph()
    p3.text = "• Establish an automated, multi-source Python ETL pipeline.\n• Quantify tail risk using Historical VaR and CVaR (95%).\n• Deliver an interactive Power BI dashboard tracking asset metrics."
    p3.level = 1

    # --- SLIDE 3: Data Architecture ---
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Data Architecture & Pipeline Flow"
    tf = slide.placeholders[1].text_frame
    tf.text = "Data Lifecycle Phases:"
    for item in [
        "1. Ingestion: Automated extraction of messy historical NAV logs and multi-variant transactions.",
        "2. Processing: Data cleaning, anomaly handling, and relational constraint enforcement via Pandas.",
        "3. Storage: Structured output exported to processed CSV databases.",
        "4. Presentation: Visual consumption via interactive business intelligence tools."
    ]:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1

    # --- SLIDE 4: ETL Operational Highlights ---
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Data Cleaning & ETL Operational Highlights"
    tf = slide.placeholders[1].text_frame
    tf.text = "Core Engineering Wins:"
    for item in [
        "• Rogue Data Containment: Mitigated structural file anomalies using selective string parsing and robust row-skipping tools (on_bad_lines='skip').",
        "• Relational Integrity: Synced isolated transaction datasets directly to core master sheets using relational AMFI codes.",
        "• Feature Engineering: Transformed static NAV records into daily percentage return metrics for downstream risk processing."
    ]:
        p = tf.add_paragraph()
        p.text = item

    # --- SLIDE 5: EDA Highlight 1 ---
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "EDA Highlight 1: Fund Performance Overview"
    tf = slide.placeholders[1].text_frame
    tf.text = "Key Structural Insights:"
    for item in [
        "• Cleansed datasets confirm significant tracking variance across 40 mutual fund schemes.",
        "• Flagship equity assets exhibit high market concentration parameters.",
        "• [Placeholder: Insert Performance Chart from Notebook]"
    ]:
        p = tf.add_paragraph()
        p.text = item

    # --- SLIDE 6: EDA Highlight 2 ---
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "EDA Highlight 2: Investor Behavioral Cohorts"
    tf = slide.placeholders[1].text_frame
    tf.text = "Behavioral Trends:"
    for item in [
        "• Onboarding Segments: Grouped users by initial transaction year to track capital retention rates.",
        "• Volume Drivers: Legacy cohorts contribute a substantial share of recurring SIP capital.",
        "• Structural Shifts: Asset preference matrices show sharp divergence between newer and long-term investors."
    ]:
        p = tf.add_paragraph()
        p.text = item

    # --- SLIDE 7: Advanced Analytics - Tail Risk ---
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Advanced Analytics: Downside Risk Quantification"
    tf = slide.placeholders[1].text_frame
    tf.text = "Risk Metric Formulations:"
    for item in [
        "• Historical VaR (95%): Establishes the 5th percentile of daily market returns, indicating the baseline boundary for expected losses.",
        "• Conditional VaR (95% CVaR): Computes the true expected shortfall by taking the average of all returns breaching the VaR threshold.",
        "• Operational Storage: Calculations run seamlessly across all active models and auto-export to reports/var_cvar_report.csv."
    ]:
        p = tf.add_paragraph()
        p.text = item

    # --- SLIDE 8: Advanced Analytics - Rolling Sharpe ---
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Advanced Analytics: 90-Day Rolling Sharpe Analysis"
    tf = slide.placeholders[1].text_frame
    tf.text = "Volatility Mapping:"
    for item in [
        "• Risk-Adjusted Tracking: Evaluated rolling consistency over time using (Rolling Mean / Rolling Std) x sqrt(252).",
        "• Core Findings: Flagship trackers show durable stability, while speculative sectors reveal sharp performance decay during market corrections.",
        "• [Placeholder: Insert rolling_sharpe_chart.png]"
    ]:
        p = tf.add_paragraph()
        p.text = item

    # --- SLIDE 9: Dashboard Showcase - Executive ---
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Power BI Dashboard Showcase: Executive View"
    tf = slide.placeholders[1].text_frame
    tf.text = "System Overview Monitoring:"
    for item in [
        "• High-level macro views mapping total Assets Under Management (AUM) and total user footprint changes.",
        "• Scalable layout optimized with clean, action-oriented KPI cards for rapid executive oversight.",
        "• [Placeholder: Insert Main Dashboard Page Screenshot]"
    ]:
        p = tf.add_paragraph()
        p.text = item

    # --- SLIDE 10: Dashboard Showcase - Deep Dive ---
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Power BI Dashboard Showcase: Granular Deep Dive"
    tf = slide.placeholders[1].text_frame
    tf.text = "Interactivity Highlights:"
    for item in [
        "• Integrated slicer panels allowing zero-latency filtering across date brackets, asset class, and risk categories.",
        "• Embedded churn detection monitors for immediate account risk evaluation.",
        "• [Placeholder: Insert Filtered Dashboard Deep Dive Screenshot]"
    ]:
        p = tf.add_paragraph()
        p.text = item

    # --- SLIDE 11: Strategic Recommendations ---
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Strategic Conclusions & Action Plan"
    tf = slide.placeholders[1].text_frame
    tf.text = "Operational Strategies:"
    for item in [
        "• Risk Mitigation: Enforce strict sector caps on highly concentrated (high HHI) profiles to mitigate volatility.",
        "• Retentive Engagement: Flag at-risk SIP accounts with payment gaps over 35 days for automated customer service outreach.",
        "• Product Engineering: Implement recommender.py's sorting logic directly into retail interfaces to drive personalized discovery."
    ]:
        p = tf.add_paragraph()
        p.text = item

    # --- SLIDE 12: Project Status & Closing ---
    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = "Architecture Status & Closing (Q&A)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Deployment Summary:"
    for item in [
        "• Production Status: Release Build v1.0 Live.",
        "• Version Control: Fully committed, clean, tagged environment successfully deployed on GitHub.",
        "• Thank you! Open for technical, architectural, or analytics questions."
    ]:
        p = tf.add_paragraph()
        p.text = item

    # Save presentation back up to root or direct directory
    output_path = "../Bluestock_MF_Presentation.pptx"
    prs.save(output_path)
    print(f"🎉 Success! Presentation generated and saved to: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_bluestock_presentation()